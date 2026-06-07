import csv
import os
from io import BytesIO
from typing import Any, Callable, Dict, Optional

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

try:  # textract is optional; gracefully degrade when it is not installed
    import textract as _textract  # type: ignore
except ModuleNotFoundError:
    _textract = None


def _read_limit(max_chars: int | None) -> int:
    return max_chars + 1 if isinstance(max_chars, int) and max_chars > 0 else -1


def _reached_limit(length: int, max_chars: int | None) -> bool:
    return isinstance(max_chars, int) and max_chars > 0 and length > max_chars


def _process_txt(file_path: str, *, max_chars: int | None = None, **_: Any) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as infile:
        return infile.read(_read_limit(max_chars))


def _process_csv(file_path: str, *, max_chars: int | None = None, **_: Any) -> str:
    rows = []
    output_length = 0
    with open(file_path, "r", encoding="utf-8-sig", errors="ignore", newline="") as infile:
        reader = csv.reader(infile)
        for row in reader:
            normalized = " ".join(cell.strip() for cell in row).rstrip()
            if normalized:
                rows.append(normalized)
                output_length += len(normalized) + 1
                if _reached_limit(output_length, max_chars):
                    break
    return "\n".join(rows)


def _process_pdf(
    file_path: str,
    *,
    page: int | None = None,
    page_from: int | None = None,
    page_to: int | None = None,
    max_chars: int | None = None,
    **_: Any,
) -> str:
    reader = PdfReader(file_path)
    total_pages = len(reader.pages)
    if total_pages == 0:
        return ""

    if page is not None:
        start_index = max(0, page - 1)
        end_index = start_index + 1
    else:
        start_index = max(0, (page_from - 1) if page_from is not None else 0)
        end_index = page_to if page_to is not None else total_pages

    end_index = min(total_pages, max(start_index + 1, end_index))
    chunks = []
    output_length = 0
    for current_page_index in range(start_index, end_index):
        current_page = reader.pages[current_page_index]
        text = current_page.extract_text() or ""
        if text:
            chunk = f"[Page {current_page_index + 1}]\n{text}"
            chunks.append(chunk)
            output_length += len(chunk) + 1
            if _reached_limit(output_length, max_chars):
                break
    return "\n".join(chunks)


def _process_docx(file_path: str, *, max_chars: int | None = None, **_: Any) -> str:
    document = Document(file_path)
    paragraphs = []
    output_length = 0
    for paragraph in document.paragraphs:
        paragraphs.append(paragraph.text)
        output_length += len(paragraph.text) + 1
        if _reached_limit(output_length, max_chars):
            break
    return "\n".join(paragraphs)


def _process_xlsx(
    file_path: str,
    *,
    sheet_name: str | None = None,
    sheet_index: int | None = None,
    max_chars: int | None = None,
    **_: Any,
) -> str:
    with open(file_path, "rb") as infile:
        workbook = load_workbook(BytesIO(infile.read()), data_only=True)
    worksheets = workbook.worksheets
    if sheet_name:
        worksheets = [sheet for sheet in worksheets if sheet.title == sheet_name]
    elif sheet_index is not None:
        if 0 <= sheet_index < len(worksheets):
            worksheets = [worksheets[sheet_index]]
        else:
            worksheets = []

    sheet_chunks = []
    output_length = 0
    for sheet in worksheets:
        rows = [f"[Sheet: {sheet.title}]"]
        for row in sheet.iter_rows(values_only=True):
            normalized = " ".join("" if cell is None else str(cell) for cell in row)
            normalized = normalized.rstrip()
            if normalized:
                rows.append(normalized)
                output_length += len(normalized) + 1
                if _reached_limit(output_length, max_chars):
                    break
        sheet_chunks.append("\n".join(rows).rstrip())
        if _reached_limit(output_length, max_chars):
            break
    workbook.close()
    return "\n\n".join(chunk for chunk in sheet_chunks if chunk)


def _shape_text(shape: Any) -> list[str]:
    chunks = []
    if getattr(shape, "has_text_frame", False):
        text = shape.text.strip()
        if text:
            chunks.append(text)
    if getattr(shape, "has_table", False):
        table = shape.table
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            normalized = " ".join(cell for cell in cells if cell).rstrip()
            if normalized:
                chunks.append(normalized)
    return chunks


def _process_pptx(file_path: str, *, max_chars: int | None = None, **_: Any) -> str:
    presentation = Presentation(file_path)
    slide_chunks = []
    output_length = 0
    for index, slide in enumerate(presentation.slides, start=1):
        chunks = [f"[Slide {index}]"]
        for shape in slide.shapes:
            chunks.extend(_shape_text(shape))

        notes_slide = getattr(slide, "notes_slide", None)
        if notes_slide:
            notes = []
            for shape in notes_slide.shapes:
                notes.extend(_shape_text(shape))
            if notes:
                chunks.append("[Notes]")
                chunks.extend(notes)

        slide_text = "\n".join(chunk for chunk in chunks if chunk).rstrip()
        if slide_text != f"[Slide {index}]":
            slide_chunks.append(slide_text)
            output_length += len(slide_text) + 2
            if _reached_limit(output_length, max_chars):
                break
    return "\n\n".join(slide_chunks)


_EXTRACTORS: Dict[str, Callable[..., str]] = {
    ".txt": _process_txt,
    ".md": _process_txt,
    ".csv": _process_csv,
    ".pdf": _process_pdf,
    ".docx": _process_docx,
    ".xlsx": _process_xlsx,
    ".pptx": _process_pptx,
}


def _extract_with_textract(file_path: str, extension: Optional[str]) -> str:
    if _textract is None:
        raise ValueError(
            "Unsupported file type without optional textract dependency. "
            "Please install textract manually if you need broader format support."
        )

    raw = _textract.process(file_path, extension=extension)
    return raw.decode("utf-8", errors="ignore")


def extract_text(
    file_path: str,
    extension: Optional[str] = None,
    *,
    page: int | None = None,
    page_from: int | None = None,
    page_to: int | None = None,
    sheet_name: str | None = None,
    sheet_index: int | None = None,
    max_chars: int | None = None,
) -> str:
    """
    Extract text content from supported files. Falls back to textract (when available)
    so that callers can still process legacy formats such as .doc.
    """
    ext = extension or os.path.splitext(file_path)[1]
    ext = (ext or "").lower()
    if not ext.startswith("."):
        ext = f".{ext}"

    extractor = _EXTRACTORS.get(ext)
    if extractor:
        return extractor(
            file_path,
            page=page,
            page_from=page_from,
            page_to=page_to,
            sheet_name=sheet_name,
            sheet_index=sheet_index,
            max_chars=max_chars,
        )

    return _extract_with_textract(file_path, ext or None)
