from __future__ import annotations

import asyncio
import os
import tempfile
import unittest

from openpyxl import Workbook
from pptx import Presentation

from app.utils.extract_text import _extract_document_in_process, extract_text_from_file
from app.utils.text_extractor import extract_text


class TextExtractorTests(unittest.TestCase):
    def test_document_process_extracts_text_successfully(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = os.path.join(temp_dir, "notes.txt")
            with open(file_path, "w", encoding="utf-8") as outfile:
                outfile.write("process extraction works")

            extracted = asyncio.run(
                extract_text_from_file(
                    file_path,
                    ".txt",
                    timeout_seconds=5,
                    max_chars=100,
                )
            )

            self.assertEqual(extracted, "process extraction works")

    def test_document_process_is_terminated_on_timeout(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = os.path.join(temp_dir, "notes.txt")
            with open(file_path, "w", encoding="utf-8") as outfile:
                outfile.write("hello")

            with self.assertRaises(TimeoutError):
                _extract_document_in_process(file_path, ".txt", {}, 0, 100)

    def test_extract_xlsx_from_path_without_extension(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            workbook_path = os.path.join(temp_dir, "source.xlsx")
            workbook = Workbook()
            sheet = workbook.active
            sheet.title = "Summary"
            sheet["A1"] = "Name"
            sheet["B1"] = "Value"
            sheet["A2"] = "Demo"
            sheet["B2"] = 42
            workbook.save(workbook_path)
            workbook.close()

            suffixless_path = os.path.join(temp_dir, "uploaded-file")
            with open(workbook_path, "rb") as source, open(suffixless_path, "wb") as target:
                target.write(source.read())

            extracted = extract_text(suffixless_path, ".xlsx")

            self.assertIn("[Sheet: Summary]", extracted)
            self.assertIn("Name Value", extracted)
            self.assertIn("Demo 42", extracted)

    def test_extract_markdown_as_text(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = os.path.join(temp_dir, "notes.md")
            with open(file_path, "w", encoding="utf-8") as outfile:
                outfile.write("# Roadmap\n\n- Add markdown upload\n")

            extracted = extract_text(file_path, ".md")

            self.assertIn("# Roadmap", extracted)
            self.assertIn("Add markdown upload", extracted)

    def test_extract_text_stops_at_character_budget(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = os.path.join(temp_dir, "large.txt")
            with open(file_path, "w", encoding="utf-8") as outfile:
                outfile.write("a" * 10_000)

            extracted = extract_text(file_path, ".txt", max_chars=100)

            self.assertEqual(len(extracted), 101)

    def test_extract_csv_rows_as_text(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = os.path.join(temp_dir, "data.csv")
            with open(file_path, "w", encoding="utf-8", newline="") as outfile:
                outfile.write("Name,Value\nDemo,42\n")

            extracted = extract_text(file_path, ".csv")

            self.assertIn("Name Value", extracted)
            self.assertIn("Demo 42", extracted)

    def test_extract_pptx_slide_text(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = os.path.join(temp_dir, "deck.pptx")
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            slide.shapes.title.text = "Quarterly Review"
            text_box = slide.shapes.add_textbox(914400, 1371600, 5486400, 914400)
            text_box.text = "Revenue grew by 12 percent"
            presentation.save(file_path)

            suffixless_path = os.path.join(temp_dir, "uploaded-deck")
            with open(file_path, "rb") as source, open(suffixless_path, "wb") as target:
                target.write(source.read())

            extracted = extract_text(suffixless_path, ".pptx")

            self.assertIn("[Slide 1]", extracted)
            self.assertIn("Quarterly Review", extracted)
            self.assertIn("Revenue grew by 12 percent", extracted)


if __name__ == "__main__":
    unittest.main()
